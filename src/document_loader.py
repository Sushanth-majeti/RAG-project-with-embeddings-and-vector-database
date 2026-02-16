"""
Document loader for multiple file formats (PDF, DOCX, XLSX, PPTX, Markdown).
"""
import os
import logging
from typing import List, Tuple, Optional
from pathlib import Path
import re

logger = logging.getLogger(__name__)

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class DocumentLoader:
    """
    Load and extract text from multiple document formats.
    Preserves document structure (headings, tables, etc.) when possible.
    """

    SUPPORTED_FORMATS = {'.pdf', '.docx', '.xlsx', '.pptx', '.md'}

    def __init__(self):
        self.supported_formats = self.SUPPORTED_FORMATS

    def load_documents(self, folder_path: str) -> List[Tuple[str, str]]:
        """
        Recursively load documents from a folder.
        Returns: List of (source_file, content) tuples
        """
        documents = []
        folder_path = Path(folder_path)

        if not folder_path.exists():
            logger.warning(f"Folder {folder_path} does not exist")
            return documents

        for file_path in folder_path.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_formats:
                try:
                    content = self._load_file(str(file_path))
                    if content.strip():
                        documents.append((str(file_path), content))
                        logger.info(f"Loaded: {file_path.name}")
                except Exception as e:
                    logger.error(f"Error loading {file_path}: {e}")

        logger.info(f"Total documents loaded: {len(documents)}")
        return documents

    def _load_file(self, file_path: str) -> str:
        """Load a single file based on its extension."""
        ext = Path(file_path).suffix.lower()

        if ext == '.pdf':
            return self._load_pdf(file_path)
        elif ext == '.docx':
            return self._load_docx(file_path)
        elif ext == '.xlsx':
            return self._load_xlsx(file_path)
        elif ext == '.pptx':
            return self._load_pptx(file_path)
        elif ext == '.md':
            return self._load_markdown(file_path)
        else:
            raise ValueError(f"Unsupported format: {ext}")

    def _load_pdf(self, file_path: str) -> str:
        """Extract text from PDF."""
        if PyPDF2 is None:
            raise ImportError("PyPDF2 is required for PDF support")

        text = []
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text.append(f"--- Page {page_num + 1} ---\n{page_text}")
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")

        return '\n'.join(text)

    def _load_docx(self, file_path: str) -> str:
        """Extract text from DOCX."""
        if DocxDocument is None:
            raise ImportError("python-docx is required for DOCX support")

        text = []
        try:
            doc = DocxDocument(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text.append(para.text)

            # Extract table content
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(
                        cell.text.strip() for cell in row.cells
                    )
                    if row_text.strip():
                        text.append(row_text)
        except Exception as e:
            logger.error(f"Error reading DOCX {file_path}: {e}")

        return '\n'.join(text)

    def _load_xlsx(self, file_path: str) -> str:
        """Extract text from XLSX."""
        if openpyxl is None:
            raise ImportError("openpyxl is required for XLSX support")

        text = []
        try:
            wb = openpyxl.load_workbook(file_path)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text.append(f"--- Sheet: {sheet} ---")
                for row in ws.iter_rows(values_only=True):
                    row_text = ' | '.join(
                        str(cell).strip() for cell in row if cell is not None
                    )
                    if row_text.strip():
                        text.append(row_text)
        except Exception as e:
            logger.error(f"Error reading XLSX {file_path}: {e}")

        return '\n'.join(text)

    def _load_pptx(self, file_path: str) -> str:
        """Extract text from PPTX."""
        if Presentation is None:
            raise ImportError("python-pptx is required for PPTX support")

        text = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                text.append(f"--- Slide {slide_num + 1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
        except Exception as e:
            logger.error(f"Error reading PPTX {file_path}: {e}")

        return '\n'.join(text)

    def _load_markdown(self, file_path: str) -> str:
        """Extract text from Markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return content
        except Exception as e:
            logger.error(f"Error reading Markdown {file_path}: {e}")
            return ""
